from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Define styles
title_font_size = Pt(44)
content_font_size = Pt(24)
blue = RGBColor(0, 51, 153)

def add_slide(title, content_lines):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1].text_frame
    for line in content_lines:
        p = content.add_paragraph()
        p.text = line
        p.font.size = content_font_size
        p.font.color.rgb = blue
    return slide

# Slide 1 - Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "The OSI Model"
slide.placeholders[1].text = "Understanding How Networks Communicate"

# Slide 2 - Introduction
add_slide("Introduction to the OSI Model", [
    "Developed by ISO to standardize network communication.",
    "Divides communication into 7 distinct layers.",
    "Helps understand how different systems interact in a network."
])

# Slide 3 - The 7 Layers Overview
add_slide("The 7 Layers of the OSI Model", [
    "7 - Application: User interface & software communication",
    "6 - Presentation: Data translation, encryption, compression",
    "5 - Session: Manages sessions between applications",
    "4 - Transport: Reliable data transfer (TCP/UDP)",
    "3 - Network: Routing and addressing (IP)",
    "2 - Data Link: Node-to-node delivery, MAC addresses",
    "1 - Physical: Hardware transmission, cables, signals"
])

# Individual layer slides
layers = [
    ("Layer 1 – Physical", [
        "Transmits raw bits over a physical medium.",
        "Examples: Cables, switches, hubs, Wi-Fi signals."
    ]),
    ("Layer 2 – Data Link", [
        "Provides node-to-node data transfer and detects errors.",
        "Examples: MAC addresses, Ethernet, switches."
    ]),
    ("Layer 3 – Network", [
        "Determines the best path for data packets (routing).",
        "Examples: IP, routers."
    ]),
    ("Layer 4 – Transport", [
        "Ensures reliable data transfer (TCP/UDP).",
        "Examples: TCP, UDP, port numbers."
    ]),
    ("Layer 5 – Session", [
        "Establishes, maintains, and terminates connections.",
        "Examples: API sessions, remote procedure calls."
    ]),
    ("Layer 6 – Presentation", [
        "Data formatting, encryption, compression.",
        "Examples: SSL/TLS, JPEG, MP3, ASCII."
    ]),
    ("Layer 7 – Application", [
        "User-level interaction with the network.",
        "Examples: HTTP, FTP, SMTP, DNS."
    ])
]
for title, content in layers:
    add_slide(title, content)

# Slide 11 - Data Encapsulation
add_slide("Data Encapsulation & Decapsulation", [
    "Data moves down the OSI layers during sending (encapsulation).",
    "Moves up the OSI layers during receiving (decapsulation).",
    "Each layer adds or removes headers for proper communication."
])

# Slide 12 - Real World Example
add_slide("Real-World Example: Web Request", [
    "User sends an HTTP request → passes through OSI layers.",
    "Server responds → data travels back through the layers.",
    "Example protocols: HTTP, TCP/IP, Ethernet, Wi-Fi."
])

# Slide 13 - Summary
add_slide("Summary", [
    "The OSI model standardizes how devices communicate.",
    "Each layer serves a unique purpose.",
    "Encapsulation, routing, encryption, and transmission all occur in layers."
])

# Slide 14 - Quiz
add_slide("Quick Review", [
    "Which layer handles encryption? (Answer: Presentation)",
    "What layer is responsible for IP addressing? (Answer: Network)",
    "Match protocols to their layers: HTTP, IP, Ethernet."
])

# Save file
prs.save("OSI_Model_Presentation.pptx")
print("✅ Presentation saved as 'OSI_Model_Presentation.pptx'")
